from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x3E)   # deep navy
ACCENT      = RGBColor(0x00, 0x8B, 0xFF)   # electric blue
LIGHT_BG    = RGBColor(0xF4, 0xF7, 0xFF)   # near-white blue tint
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GOLD        = RGBColor(0xFF, 0xBF, 0x00)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
MID_GREY    = RGBColor(0x5A, 0x6A, 0x85)
CARD_BG     = RGBColor(0xE8, 0xF1, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helper utilities ────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb, transparency=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.line.width = 0
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_rgb
    return shape


def add_text_box(slide, text, l, t, w, h,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_bullet_box(slide, items, l, t, w, h,
                   font_size=16, color=DARK_TEXT,
                   bold_first_word=False, line_spacing=1.2):
    """items = list of (icon_str, text) tuples"""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (icon, text) in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_before = Pt(4)
        p.space_after  = Pt(4)

        # icon run
        r_icon = p.add_run()
        r_icon.text = icon + "  "
        r_icon.font.size  = Pt(font_size)
        r_icon.font.color.rgb = ACCENT
        r_icon.font.name  = "Segoe UI Emoji"

        # text run
        r_text = p.add_run()
        r_text.text = text
        r_text.font.size  = Pt(font_size)
        r_text.font.color.rgb = color
        r_text.font.name  = "Calibri"
    return txBox


def add_accent_line(slide, l, t, w, thickness=0.04):
    add_rect(slide, l, t, w, thickness, ACCENT)


def add_card(slide, l, t, w, h, icon, title, body, title_color=NAVY, body_color=MID_GREY):
    """A rounded-ish card: background + icon + title + body."""
    add_rect(slide, l, t, w, h, CARD_BG)
    # top accent strip
    add_rect(slide, l, t, w, 0.05, ACCENT)
    # icon
    add_text_box(slide, icon, l+0.15, t+0.1, 0.6, 0.55,
                 font_size=26, color=ACCENT, bold=False)
    # title
    add_text_box(slide, title, l+0.15, t+0.6, w-0.3, 0.4,
                 font_size=13, bold=True, color=title_color)
    # body
    add_text_box(slide, body, l+0.15, t+0.95, w-0.3, h-1.1,
                 font_size=11, color=body_color, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)

# Full navy background
add_rect(s1, 0, 0, 13.33, 7.5, NAVY)

# Left electric-blue sidebar
add_rect(s1, 0, 0, 0.55, 7.5, ACCENT)

# Gold accent bar under title area
add_rect(s1, 0.55, 3.85, 12.78, 0.07, GOLD)

# Top-right decorative circle (subtle)
add_rect(s1, 10.5, -1.2, 4.0, 4.0, RGBColor(0x0A, 0x29, 0x5C))

# Project label chip
add_rect(s1, 1.0, 1.0, 3.2, 0.45, ACCENT)
add_text_box(s1, "AI  PRODUCT  SHOWCASE", 1.05, 1.02, 3.1, 0.4,
             font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Main title
add_text_box(s1, "FMCG Deal Intelligence", 1.0, 1.75, 10.5, 1.1,
             font_size=46, bold=True, color=WHITE)
add_text_box(s1, "Newsletter Platform", 1.0, 2.75, 10.5, 0.9,
             font_size=46, bold=True, color=ACCENT)

# Tagline
add_text_box(s1,
             "Turning unstructured market noise into actionable FMCG intelligence — automatically.",
             1.0, 4.05, 9.5, 0.7, font_size=16, italic=True,
             color=RGBColor(0xB0, 0xC8, 0xFF))

# Divider
add_rect(s1, 1.0, 4.95, 5.5, 0.03, MID_GREY)

# Name & role
add_text_box(s1, "Akshat Banga", 1.0, 5.1, 5.0, 0.5,
             font_size=18, bold=True, color=WHITE)
add_text_box(s1, "AI / ML Engineer  ·  Full-Stack AI Developer", 1.0, 5.6, 7.0, 0.4,
             font_size=12, color=RGBColor(0x8A, 0xA8, 0xD8))

# Date bottom-right
add_text_box(s1, "March 2026", 10.5, 6.8, 2.5, 0.4,
             font_size=11, color=MID_GREY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
add_rect(s2, 0, 0, 13.33, 7.5, LIGHT_BG)
add_rect(s2, 0, 0, 13.33, 1.5, NAVY)
add_rect(s2, 0, 0, 0.55, 7.5, ACCENT)

# Slide header
add_text_box(s2, "02", 0.7, 0.15, 0.9, 0.6,
             font_size=28, bold=True, color=RGBColor(0x1F, 0x4A, 0x9C))
add_text_box(s2, "Problem Statement", 1.5, 0.18, 8.0, 0.6,
             font_size=26, bold=True, color=WHITE)
add_text_box(s2, "The FMCG intelligence gap that slows every analyst down",
             1.5, 0.75, 10.0, 0.5, font_size=13, italic=True,
             color=RGBColor(0x8A, 0xB4, 0xE8))

# Problem cards  (3 cards across)
cards = [
    ("🌊", "Information Overload",
     "100s of news articles, press releases & market reports published daily.\nNo single source captures the full picture."),
    ("⏱️", "Manual Research Drain",
     "Analysts spend 3–5 hours per day just gathering & reading news.\nHigh effort, low output — not scalable."),
    ("🧩", "No Structured Insights",
     "Raw news lacks context: deals aren't tagged, trends aren't surfaced, signals are buried in text."),
]
positions = [0.75, 4.75, 8.75]
for (icon, title, body), x in zip(cards, positions):
    add_card(s2, x, 1.9, 3.7, 3.8, icon, title, body)

# Bottom callout
add_rect(s2, 0.75, 5.95, 11.8, 0.9, NAVY)
add_text_box(s2,
             "💡  The result: slow, fragmented, and inconsistent market intelligence — "
             "when speed and accuracy matter most.",
             1.0, 6.05, 11.4, 0.7, font_size=13, color=WHITE, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution Overview
# ════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
add_rect(s3, 0, 0, 13.33, 7.5, LIGHT_BG)
add_rect(s3, 0, 0, 13.33, 1.5, NAVY)
add_rect(s3, 0, 0, 0.55, 7.5, ACCENT)

add_text_box(s3, "03", 0.7, 0.15, 0.9, 0.6,
             font_size=28, bold=True, color=RGBColor(0x1F, 0x4A, 0x9C))
add_text_box(s3, "Solution Overview", 1.5, 0.18, 8.0, 0.6,
             font_size=26, bold=True, color=WHITE)
add_text_box(s3, "An AI-powered platform that does the heavy lifting — end to end",
             1.5, 0.75, 10.0, 0.5, font_size=13, italic=True,
             color=RGBColor(0x8A, 0xB4, 0xE8))

# Hero statement box
add_rect(s3, 0.75, 1.75, 11.8, 1.0, ACCENT)
add_text_box(s3,
             "\"  Turning unstructured FMCG news into structured, AI-curated intelligence — "
             "delivered as an interactive newsletter dashboard.  \"",
             1.0, 1.85, 11.4, 0.8, font_size=15, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, italic=True)

# 4 solution bullets
bullets = [
    ("🔍", "Automated News Aggregation — continuously pulls FMCG articles, press releases & deal announcements"),
    ("🤖", "LLM-Powered Summarisation — AI condenses lengthy articles into crisp, 2–3 line executive summaries"),
    ("🏷️", "Smart Categorisation — deals auto-tagged by type: Funding, M&A, Product Launch, Expansion & more"),
    ("📊", "Interactive Dashboard — real-time, filterable frontend for instant access to curated intelligence"),
]
add_bullet_box(s3, bullets, 0.85, 2.95, 11.6, 3.8,
               font_size=15, color=DARK_TEXT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — System / AI Architecture
# ════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
add_rect(s4, 0, 0, 13.33, 7.5, LIGHT_BG)
add_rect(s4, 0, 0, 13.33, 1.5, NAVY)
add_rect(s4, 0, 0, 0.55, 7.5, ACCENT)

add_text_box(s4, "04", 0.7, 0.15, 0.9, 0.6,
             font_size=28, bold=True, color=RGBColor(0x1F, 0x4A, 0x9C))
add_text_box(s4, "System Architecture", 1.5, 0.18, 8.0, 0.6,
             font_size=26, bold=True, color=WHITE)
add_text_box(s4, "A clean 4-layer pipeline: Ingest → Process → Analyse → Deliver",
             1.5, 0.75, 10.0, 0.5, font_size=13, italic=True,
             color=RGBColor(0x8A, 0xB4, 0xE8))

# Pipeline stages — 4 boxes with arrows
stages = [
    ("🌐", "DATA\nINGESTION", "News APIs\nPress Releases\nRSS Feeds\nWeb Crawlers"),
    ("⚙️", "PROCESSING\nLAYER",  "Text Cleaning\nDe-duplication\nParsing\nMetadata Tagging"),
    ("🧠", "AI /\nLLM ENGINE",   "Claude / GPT\nSummarisation\nEntity Extraction\nSentiment Scoring"),
    ("🖥️", "FRONTEND\nDELIVERY", "Next.js Dashboard\nFilterable Feed\nNewsletter View\nReal-time Updates"),
]

box_w, box_h = 2.6, 3.5
start_x = 0.8
gap = 0.3
arrow_y = 3.35

for i, (icon, title, body) in enumerate(stages):
    x = start_x + i * (box_w + gap + 0.35)
    # Card bg
    add_rect(s4, x, 1.85, box_w, box_h, NAVY)
    add_rect(s4, x, 1.85, box_w, 0.06, ACCENT)
    # Icon
    add_text_box(s4, icon, x + 0.9, 2.0, 0.9, 0.6, font_size=28, color=WHITE)
    # Title
    add_text_box(s4, title, x + 0.1, 2.65, box_w - 0.2, 0.75,
                 font_size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    # Body
    add_text_box(s4, body, x + 0.15, 3.45, box_w - 0.3, 1.7,
                 font_size=11, color=RGBColor(0xB0, 0xC8, 0xFF),
                 align=PP_ALIGN.CENTER)
    # Arrow (except after last)
    if i < 3:
        ax = x + box_w + 0.05
        add_text_box(s4, "→", ax, arrow_y - 0.05, 0.4, 0.5,
                     font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Tech stack strip
add_rect(s4, 0.75, 5.6, 11.8, 0.65, CARD_BG)
add_text_box(s4, "Tech Stack:", 0.95, 5.67, 1.5, 0.45,
             font_size=12, bold=True, color=NAVY)
add_text_box(s4,
             "Python · FastAPI · LLM APIs (Claude/GPT) · Next.js · TypeScript · "
             "Railway (Backend) · Vercel (Frontend)",
             2.35, 5.67, 10.0, 0.45, font_size=12, color=MID_GREY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Key Features
# ════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
add_rect(s5, 0, 0, 13.33, 7.5, LIGHT_BG)
add_rect(s5, 0, 0, 13.33, 1.5, NAVY)
add_rect(s5, 0, 0, 0.55, 7.5, ACCENT)

add_text_box(s5, "05", 0.7, 0.15, 0.9, 0.6,
             font_size=28, bold=True, color=RGBColor(0x1F, 0x4A, 0x9C))
add_text_box(s5, "Key Features", 1.5, 0.18, 8.0, 0.6,
             font_size=26, bold=True, color=WHITE)
add_text_box(s5, "Built for analysts who need intelligence, not just information",
             1.5, 0.75, 10.0, 0.5, font_size=13, italic=True,
             color=RGBColor(0x8A, 0xB4, 0xE8))

# 6 feature cards, 2 rows × 3 cols
features = [
    ("🤖", "Smart AI Summarisation",
     "LLM reduces full articles to sharp 2–3 line summaries without losing key context."),
    ("💼", "Deal Intelligence Extraction",
     "Automatically identifies & structures M&A deals, funding rounds, and partnerships."),
    ("🏷️", "Auto-Categorisation",
     "Tags every story: Funding · M&A · Expansion · Product · Regulatory."),
    ("📱", "Interactive Dashboard",
     "Futuristic Next.js UI with real-time filters, search & responsive design."),
    ("⚡", "Workflow Automation",
     "Zero manual curation — pipeline runs end-to-end with no human-in-the-loop."),
    ("🔌", "API-First Architecture",
     "FastAPI backend exposes clean endpoints; easily integrates with any consumer."),
]

cols, rows = 3, 2
fw, fh = 3.7, 2.3
col_gap = 0.38
row_gap = 0.25
start_fx, start_fy = 0.75, 1.75

for idx, (icon, title, body) in enumerate(features):
    col = idx % cols
    row = idx // cols
    x = start_fx + col * (fw + col_gap)
    y = start_fy + row * (fh + row_gap)
    add_rect(s5, x, y, fw, fh, NAVY)
    add_rect(s5, x, y, fw, 0.05, ACCENT)
    add_text_box(s5, icon, x + 0.15, y + 0.1, 0.6, 0.55, font_size=24, color=ACCENT)
    add_text_box(s5, title, x + 0.75, y + 0.12, fw - 0.9, 0.45,
                 font_size=13, bold=True, color=WHITE)
    add_text_box(s5, body, x + 0.15, y + 0.65, fw - 0.3, 1.45,
                 font_size=11, color=RGBColor(0xB0, 0xC8, 0xFF), wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Results & Impact
# ════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
add_rect(s6, 0, 0, 13.33, 7.5, LIGHT_BG)
add_rect(s6, 0, 0, 13.33, 1.5, NAVY)
add_rect(s6, 0, 0, 0.55, 7.5, ACCENT)

add_text_box(s6, "06", 0.7, 0.15, 0.9, 0.6,
             font_size=28, bold=True, color=RGBColor(0x1F, 0x4A, 0x9C))
add_text_box(s6, "Results & Business Impact", 1.5, 0.18, 9.0, 0.6,
             font_size=26, bold=True, color=WHITE)
add_text_box(s6, "Measurable value — from day one",
             1.5, 0.75, 10.0, 0.5, font_size=13, italic=True,
             color=RGBColor(0x8A, 0xB4, 0xE8))

# KPI metric boxes
metrics = [
    ("~70%", "Reduction in manual\nresearch effort (est.)"),
    ("5×",   "Faster access to\nmarket insights"),
    ("100%", "Automation of the\ncuration workflow"),
    ("∞",    "Scalable — handles\nany volume of news"),
]
mw, mh = 2.7, 1.9
mx_start = 0.75
for i, (val, label) in enumerate(metrics):
    mx = mx_start + i * (mw + 0.4)
    add_rect(s6, mx, 1.75, mw, mh, NAVY)
    add_rect(s6, mx, 1.75, mw, 0.06, GOLD)
    add_text_box(s6, val, mx, 1.85, mw, 0.85,
                 font_size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(s6, label, mx + 0.1, 2.7, mw - 0.2, 0.85,
                 font_size=11, color=RGBColor(0xB0, 0xC8, 0xFF),
                 align=PP_ALIGN.CENTER)

# Impact bullets
impacts = [
    ("✅", "Analysts get a structured daily briefing — no more tab-hopping across 10 sites"),
    ("✅", "Decision-makers spot emerging FMCG trends before competitors do"),
    ("✅", "Consistent, bias-free AI summaries replace inconsistent manual notes"),
    ("✅", "Platform scales to cover new categories, regions & data sources with zero rework"),
]
add_bullet_box(s6, impacts, 0.85, 3.85, 11.6, 3.2,
               font_size=15, color=DARK_TEXT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Future Scope / Vision
# ════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
add_rect(s7, 0, 0, 13.33, 7.5, NAVY)
add_rect(s7, 0, 0, 0.55, 7.5, ACCENT)
# subtle top-right circle
add_rect(s7, 9.5, -0.5, 4.5, 4.5, RGBColor(0x0A, 0x29, 0x5C))

add_text_box(s7, "07", 0.7, 0.15, 0.9, 0.6,
             font_size=28, bold=True, color=RGBColor(0x1F, 0x6A, 0xDC))
add_text_box(s7, "Future Scope & Vision", 1.5, 0.18, 9.0, 0.6,
             font_size=26, bold=True, color=WHITE)
add_text_box(s7, "From newsletter to AI-powered decision intelligence platform",
             1.5, 0.75, 10.5, 0.5, font_size=13, italic=True,
             color=RGBColor(0x8A, 0xB4, 0xE8))

add_accent_line(s7, 0.75, 1.5, 11.8)

future_items = [
    ("🔔", "Real-Time Alerts — push notifications when a significant deal breaks in tracked categories"),
    ("📈", "Predictive Analytics — ML models forecasting deal likelihood based on company signals"),
    ("🤝", "AI Decision Agents — autonomous agents that recommend actions (invest / monitor / avoid)"),
    ("🌍", "Multi-Market Expansion — extend beyond India to SEA, Middle East & global FMCG markets"),
    ("🔗", "CRM & Workflow Integration — plug directly into Salesforce, Notion, or internal portals"),
]
add_bullet_box(s7, future_items, 0.85, 1.7, 11.6, 4.2,
               font_size=15, color=WHITE)

# Closing banner
add_rect(s7, 0.75, 6.15, 11.8, 1.0, ACCENT)
add_text_box(s7,
             "The foundation is built.  The intelligence layer is live.  "
             "The next step: a fully autonomous FMCG market AI.",
             1.0, 6.27, 11.4, 0.75,
             font_size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, italic=True)


# ── Save ─────────────────────────────────────────────────────────────────────
output_path = "/home/user/FMCG-deal-intelligence-newsletter/FMCG_Deal_Intelligence_Presentation.pptx"
prs.save(output_path)
print(f"Saved → {output_path}")
